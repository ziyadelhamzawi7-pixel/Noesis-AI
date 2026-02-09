"""
PowerPoint parsing tool with text extraction from slides, notes, and shapes.
Supports .pptx files using python-pptx library.
"""

import sys
from pathlib import Path
from typing import Dict, List, Any, Optional
import json
from loguru import logger

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    logger.error("Required package not installed. Run: pip install python-pptx")
    sys.exit(1)


class PowerPointParser:
    """PowerPoint parsing with comprehensive text extraction."""

    def __init__(self, file_path: Path):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")

    def parse(self) -> Dict[str, Any]:
        """
        Parse PowerPoint file and extract all text content.

        Returns:
            Dictionary with extracted content including slides, notes, and metadata
        """
        logger.info(f"Parsing PowerPoint: {self.file_path.name}")

        result = {
            "file_name": self.file_path.name,
            "file_path": str(self.file_path.absolute()),
            "file_size": self.file_path.stat().st_size,
            "slides": [],
            "text": "",
            "tables": [],
            "metadata": {},
            "method": "python-pptx",
            "slide_count": 0,
            "has_notes": False,
            "has_tables": False
        }

        # Pre-validation
        try:
            if self.file_path.stat().st_size == 0:
                result["error"] = "PowerPoint file is empty (0 bytes)"
                logger.error(f"PowerPoint file is empty: {self.file_path}")
                return result

            # Check for valid PPTX (ZIP-based format)
            with open(self.file_path, 'rb') as f:
                header = f.read(4)
                if header != b'PK\x03\x04':
                    result["error"] = f"File is not a valid PowerPoint file (header: {header!r})"
                    logger.error(f"Invalid PPTX header for {self.file_path}: {header!r}")
                    return result
        except Exception as e:
            result["error"] = f"Cannot read PowerPoint file: {str(e)}"
            logger.error(f"Cannot read PowerPoint file {self.file_path}: {e}")
            return result

        try:
            prs = Presentation(str(self.file_path))
            result["slide_count"] = len(prs.slides)

            # Extract metadata
            result["metadata"] = self._extract_metadata(prs)

            all_text = []
            all_tables = []

            # Process each slide
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_data = self._extract_slide_content(slide, slide_num)
                result["slides"].append(slide_data)

                # Collect text
                if slide_data.get("text"):
                    all_text.append(f"--- Slide {slide_num} ---\n{slide_data['text']}")

                # Collect notes
                if slide_data.get("notes"):
                    all_text.append(f"[Notes for Slide {slide_num}]\n{slide_data['notes']}")
                    result["has_notes"] = True

                # Collect tables
                if slide_data.get("tables"):
                    for table in slide_data["tables"]:
                        table["slide_number"] = slide_num
                        all_tables.append(table)
                    result["has_tables"] = True

            result["text"] = "\n\n".join(all_text)
            result["tables"] = all_tables

            # Calculate statistics
            result["total_chars"] = len(result["text"])
            result["total_tables"] = len(result["tables"])

            logger.success(
                f"Parsed {result['slide_count']} slides, "
                f"{result['total_chars']} chars, "
                f"{result['total_tables']} tables"
            )

            return result

        except Exception as e:
            logger.error(f"Unexpected error parsing PowerPoint: {e}")
            result["error"] = str(e)
            return result

    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """Extract PowerPoint metadata."""
        try:
            core_props = prs.core_properties
            return {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "comments": core_props.comments or "",
                "category": core_props.category or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "last_modified_by": core_props.last_modified_by or "",
            }
        except Exception as e:
            logger.warning(f"Failed to extract metadata: {e}")
            return {}

    def _extract_slide_content(self, slide, slide_num: int) -> Dict[str, Any]:
        """Extract all content from a single slide."""
        slide_data = {
            "slide_number": slide_num,
            "text": "",
            "notes": "",
            "shapes": [],
            "tables": [],
            "char_count": 0
        }

        text_parts = []
        tables = []

        # Extract text from all shapes
        for shape in slide.shapes:
            try:
                shape_info = self._extract_shape_content(shape)

                if shape_info.get("text"):
                    text_parts.append(shape_info["text"])
                    slide_data["shapes"].append({
                        "type": shape_info.get("type", "unknown"),
                        "name": shape.name,
                        "text_preview": shape_info["text"][:100] if shape_info["text"] else ""
                    })

                if shape_info.get("table"):
                    tables.append(shape_info["table"])

            except Exception as e:
                logger.debug(f"Failed to extract shape {shape.name}: {e}")

        # Extract speaker notes
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                if notes_frame and notes_frame.text:
                    slide_data["notes"] = notes_frame.text.strip()
        except Exception as e:
            logger.debug(f"Failed to extract notes for slide {slide_num}: {e}")

        slide_data["text"] = "\n".join(text_parts)
        slide_data["tables"] = tables
        slide_data["char_count"] = len(slide_data["text"])

        return slide_data

    def _extract_shape_content(self, shape) -> Dict[str, Any]:
        """Extract content from a shape."""
        result = {
            "type": "unknown",
            "text": "",
            "table": None
        }

        # Determine shape type
        if shape.has_text_frame:
            result["type"] = "text"
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    text_parts.append(para_text)
            result["text"] = "\n".join(text_parts)

        elif shape.has_table:
            result["type"] = "table"
            table = shape.table
            table_data = []

            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_data.append(cell_text)
                table_data.append(row_data)

            result["table"] = {
                "rows": len(table.rows),
                "columns": len(table.columns),
                "data": table_data,
                "headers": table_data[0] if table_data else []
            }

            # Also add table content as text
            table_text = self._table_to_text(table_data)
            result["text"] = table_text

        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result["type"] = "group"
            # Recursively extract from grouped shapes
            group_texts = []
            for sub_shape in shape.shapes:
                sub_result = self._extract_shape_content(sub_shape)
                if sub_result.get("text"):
                    group_texts.append(sub_result["text"])
            result["text"] = "\n".join(group_texts)

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            result["type"] = "picture"
            # Pictures don't have text, but we note their presence

        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            result["type"] = "chart"
            # Charts may have title text
            try:
                if hasattr(shape, 'chart') and shape.chart.has_title:
                    result["text"] = f"[Chart: {shape.chart.chart_title.text_frame.text}]"
            except Exception:
                result["text"] = "[Chart]"

        return result

    def _table_to_text(self, table_data: List[List[str]]) -> str:
        """Convert table data to readable text format."""
        if not table_data:
            return ""

        lines = []
        for row in table_data:
            row_text = " | ".join(cell for cell in row if cell)
            if row_text:
                lines.append(row_text)

        return "\n".join(lines)


def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    Convenience function to parse a PowerPoint file.

    Args:
        file_path: Path to PowerPoint file (.pptx)

    Returns:
        Dictionary with parsed content

    Example:
        >>> result = parse_pptx("pitch_deck.pptx")
        >>> print(f"Extracted {len(result['text'])} characters from {result['slide_count']} slides")
    """
    parser = PowerPointParser(file_path)
    return parser.parse()


def main():
    """CLI interface for testing."""
    import argparse

    parser = argparse.ArgumentParser(description="Parse PowerPoint documents")
    parser.add_argument("file", help="Path to PowerPoint file (.pptx)")
    parser.add_argument("--output", help="Output JSON file path")
    parser.add_argument("--verbose", action="store_true", help="Verbose logging")

    args = parser.parse_args()

    if args.verbose:
        logger.remove()
        logger.add(sys.stderr, level="DEBUG")

    # Parse PowerPoint
    result = parse_pptx(args.file)

    # Output
    if args.output:
        output_path = Path(args.output)
        output_path.write_text(json.dumps(result, indent=2, default=str))
        logger.info(f"Saved results to {output_path}")
    else:
        # Print summary
        print(f"\n{'='*60}")
        print(f"File: {result['file_name']}")
        print(f"Slides: {result['slide_count']}")
        print(f"Characters: {result.get('total_chars', 0):,}")
        print(f"Tables: {result.get('total_tables', 0)}")
        print(f"Has Notes: {result.get('has_notes', False)}")
        print(f"Method: {result['method']}")

        if result.get('error'):
            print(f"Error: {result['error']}")

        print(f"{'='*60}\n")

        # Print first 500 chars
        if result['text']:
            print("First 500 characters:")
            print(result['text'][:500])
            print("...")


if __name__ == "__main__":
    main()
